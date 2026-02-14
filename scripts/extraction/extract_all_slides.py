# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
import json

sys.stdout.reconfigure(encoding='utf-8')

def extract_all_slides():
    pptx_path = r"K:\문서화\보험개발원\data\보험개발원 제안서_v0.5_재정정보원 제안서 초안.pptx"
    
    try:
        prs = Presentation(pptx_path)
        print(f"=== 전체 {len(prs.slides)}개 슬라이드 분석 ===\n")
        
        all_slides = []
        
        # 모든 슬라이드 텍스트 추출
        for i, slide in enumerate(prs.slides):
            slide_content = {
                'number': i + 1,
                'texts': []
            }
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content['texts'].append(shape.text.strip())
            
            all_slides.append(slide_content)
            
            # 주요 슬라이드 상세 출력
            if i + 1 in [1, 2, 3, 4, 5, 6, 14, 15, 16, 17, 18, 19, 20, 21, 37, 38, 39, 40, 41, 42, 55, 56, 57, 58, 59, 60, 61, 62]:
                print(f"\n[슬라이드 {i+1}] {'='*50}")
                for text in slide_content['texts']:
                    if len(text) > 200:
                        print(f"• {text[:200]}...")
                    else:
                        print(f"• {text}")
        
        # 슬라이드 구성 요약
        print("\n" + "="*60)
        print("제안서 구조 분석")
        print("="*60)
        
        sections = {
            "Ⅰ. 제안개요": [],
            "Ⅱ. 기술부문": [], 
            "Ⅲ. 사업관리부문": [],
            "Ⅳ. 지원부문": [],
            "Ⅴ. 가격부문": []
        }
        
        current_section = ""
        for slide in all_slides:
            for text in slide['texts']:
                for section in sections.keys():
                    if section in text:
                        current_section = section
                        break
            if current_section:
                sections[current_section].append(slide['number'])
        
        for section, pages in sections.items():
            if pages:
                print(f"{section}: 슬라이드 {pages[0]}-{pages[-1]} ({len(pages)}장)")
        
    except Exception as e:
        print(f"오류: {e}")

if __name__ == "__main__":
    extract_all_slides()