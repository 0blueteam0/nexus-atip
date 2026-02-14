# -*- coding: utf-8 -*-
import os
import sys
from pptx import Presentation
import pypdf

# UTF-8 인코딩 설정
sys.stdout.reconfigure(encoding='utf-8')

def analyze_pptx():
    """PowerPoint 제안서 분석"""
    pptx_path = r"K:\문서화\보험개발원\data\보험개발원 제안서_v0.5_재정정보원 제안서 초안.pptx"
    
    print("="*60)
    print("보험개발원 제안서 분석")
    print("="*60)
    
    try:
        prs = Presentation(pptx_path)
        print(f"총 슬라이드 수: {len(prs.slides)}\n")
        
        # 중요 슬라이드 분석
        important_slides = [4, 17, 18]  # 0-indexed (슬라이드 5, 18, 19)
        
        for idx in important_slides:
            if idx < len(prs.slides):
                slide = prs.slides[idx]
                print(f"\n[슬라이드 {idx+1}]")
                print("-"*40)
                
                # 슬라이드의 모든 텍스트 추출
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    for text in slide_text:
                        print(f"• {text}")
                else:
                    print("(텍스트 없음)")
        
        # 전체 슬라이드 제목 목록
        print("\n" + "="*60)
        print("전체 슬라이드 구성")
        print("="*60)
        
        for i, slide in enumerate(prs.slides):
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip().split('\n')[0][:50]
                    break
            print(f"슬라이드 {i+1}: {title}")
            
    except Exception as e:
        print(f"PowerPoint 분석 오류: {e}")

def analyze_pdfs():
    """PDF 문서 분석"""
    pdf_files = [
        r"K:\문서화\보험개발원\data\붙임1. 입찰공고문_24시간 상시 통합보안관제.pdf",
        r"K:\문서화\보험개발원\data\붙임2. 제안요청서_24시간 상시 통합보안관제.pdf"
    ]
    
    for pdf_path in pdf_files:
        try:
            print("\n" + "="*60)
            print(f"분석 중: {os.path.basename(pdf_path)}")
            print("="*60)
            
            with open(pdf_path, 'rb') as file:
                pdf = pypdf.PdfReader(file)
                print(f"총 페이지 수: {len(pdf.pages)}")
                
                # 처음 3페이지 내용 추출
                for i in range(min(3, len(pdf.pages))):
                    page = pdf.pages[i]
                    text = page.extract_text()
                    if text.strip():
                        print(f"\n[페이지 {i+1}]")
                        print(text[:500])  # 처음 500자만
                        
        except Exception as e:
            print(f"PDF 분석 오류: {e}")

if __name__ == "__main__":
    analyze_pptx()
    analyze_pdfs()