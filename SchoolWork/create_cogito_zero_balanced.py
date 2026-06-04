from pathlib import Path
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.shared import Inches, Pt, RGBColor
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.util import Inches as PptInches, Pt as PptPt

BASE = Path('J:/PortableApps/genai/SchoolWork')
DOCX_OUT = BASE / 'Cogito_Zero_사업계획서_v1_1_balanced_business_plan.docx'
PPTX_OUT = BASE / 'Cogito_Zero_발표자료_v1_1_balanced_business_plan.pptx'

C = {
    'bg': 'F7F5F0',
    'charcoal': '1F2428',
    'slate': '6E7781',
    'teal': '0F766E',
    'line': 'D8D5CE',
    'white': 'FFFFFF',
    'soft': 'EFECE5',
}


def rgb(hex_value):
    h = hex_value.strip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def drgb(name):
    return RGBColor(*rgb(C[name]))


def prgb(name):
    return PptRGBColor(*rgb(C[name]))


def set_run(run, size=10, color='charcoal', bold=False):
    run.font.name = '맑은 고딕'
    run._element.rPr.rFonts.set(qn('w:eastAsia'), '맑은 고딕')
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = drgb(color)


def shade(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:fill'), fill)
    tc_pr.append(shd)


def cell_text(cell, text, bold=False, color='charcoal', size=8.6):
    cell.text = ''
    p = cell.paragraphs[0]
    r = p.add_run(text)
    set_run(r, size=size, color=color, bold=bold)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def h(doc, text, level=1):
    p = doc.add_paragraph()
    p.style = f'Heading {level}'
    r = p.add_run(text)
    set_run(r, size=17 if level == 1 else 12.5, color='charcoal', bold=True)
    p.paragraph_format.space_before = Pt(10 if level == 1 else 5)
    p.paragraph_format.space_after = Pt(5)
    return p


def para(doc, text, lead=None, size=9.8):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(5)
    p.paragraph_format.line_spacing = 1.13
    if lead and text.startswith(lead):
        r = p.add_run(lead)
        set_run(r, size=size, color='teal', bold=True)
        r2 = p.add_run(text[len(lead):])
        set_run(r2, size=size, color='charcoal')
    else:
        r = p.add_run(text)
        set_run(r, size=size, color='charcoal')
    return p


def bullets(doc, items):
    for item in items:
        p = doc.add_paragraph(style='List Bullet')
        p.paragraph_format.space_after = Pt(2)
        r = p.add_run(item)
        set_run(r, size=9.3, color='charcoal')


def table(doc, headers, rows, widths=None):
    t = doc.add_table(rows=1, cols=len(headers))
    t.style = 'Table Grid'
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, head in enumerate(headers):
        shade(t.rows[0].cells[i], C['charcoal'])
        cell_text(t.rows[0].cells[i], head, bold=True, color='white', size=8.3)
    for ri, row in enumerate(rows):
        cells = t.add_row().cells
        for i, val in enumerate(row):
            shade(cells[i], 'FFFFFF' if ri % 2 == 0 else C['soft'])
            cell_text(cells[i], val, size=8.2)
    if widths:
        for row in t.rows:
            for i, w in enumerate(widths):
                row.cells[i].width = Inches(w)
    doc.add_paragraph()
    return t


def build_docx():
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.68)
    sec.bottom_margin = Inches(0.68)
    sec.left_margin = Inches(0.74)
    sec.right_margin = Inches(0.74)
    doc.styles['Normal'].font.name = '맑은 고딕'
    doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), '맑은 고딕')

    p = doc.add_paragraph()
    r = p.add_run('Cogito Zero')
    set_run(r, size=28, bold=True)
    p2 = doc.add_paragraph()
    r = p2.add_run('AI 에이전트 실행 통제를 위한 제로트러스트 게이트웨이')
    set_run(r, size=13, color='teal', bold=True)
    p3 = doc.add_paragraph()
    r = p3.add_run('Balanced Business Plan Copy · v1.1 · 2026')
    set_run(r, size=8.8, color='slate')
    para(doc, '본 문서는 원본 상세본과 light 버전의 중간 지점을 목표로 작성한 balanced 사본입니다. 원본의 시장·거버넌스·재무·리스크 설득 요소를 일부 되살리되, 장황한 서사와 과도한 세부 설명은 줄였습니다.')
    para(doc, '핵심 정의: Cogito Zero는 AI 에이전트와 도구/MCP/API 사이에서 모든 실행 요청을 정책 기반으로 평가하고, 필요한 경우 조건부 실행·승인·드라이런·차단으로 전환하는 제로트러스트 실행 통제 게이트웨이입니다.', '핵심 정의:')

    doc.add_page_break()
    h(doc, '목차')
    toc = ['1. 프롤로그: 병목은 모델이 아니라 실행 통제다', '2. 핵심 요약', '3. 시장 기회: Agent Execution Security의 등장', '4. 문제 정의: Tool Call의 블랙박스', '5. 솔루션: Zero Trust Gateway for AI Agents', '6. 제품 구성과 작동 방식', '7. 정책 결정 체계', '8. 대표 시나리오 3가지', '9. 거버넌스와 컴플라이언스', '10. 경쟁 환경과 포지셔닝', '11. 비즈니스 모델', '12. GTM 및 초기 도입 전략', '13. 로드맵과 팀', '14. 재무 관점과 리스크', '15. 맺음말']
    bullets(doc, toc)

    doc.add_page_break()
    sections = [
        ('1. 프롤로그: 병목은 모델이 아니라 실행 통제다', [
            'AI 에이전트는 더 이상 답변을 생성하는 데서 멈추지 않습니다. MCP, API, SaaS, 데이터베이스, 보안장비를 호출하며 실제 업무를 수행하는 운영 주체로 이동하고 있습니다.',
            '이 변화가 커질수록 기업의 질문도 바뀝니다. “어떤 모델을 쓸 것인가”보다 “이 에이전트에게 어디까지 실행 권한을 줄 수 있는가”가 더 중요한 병목이 됩니다.',
            '사람에게는 계정, 권한, 승인, 감사라는 통제 체계가 있습니다. 그러나 사람 중심 체계만으로는 에이전트의 빠르고 반복적인 Tool Call을 충분히 설명하거나 통제하기 어렵습니다.',
            'Cogito Zero는 바로 이 지점에 위치합니다. 에이전트와 도구 사이에서 실행 요청을 검증하고, 위험한 호출은 조건부 실행·승인·드라이런·차단으로 전환합니다.',
        ]),
        ('2. 핵심 요약', [
            'Cogito Zero는 AI 에이전트 실행 보안을 위한 제로트러스트 게이트웨이입니다. 모든 Tool Call을 실행 직전에 평가하고, 정책에 따라 허용 또는 제한합니다.',
            '제품의 핵심 가치는 자동화를 멈추지 않으면서 위험한 실행만 선별적으로 제어하는 데 있습니다. 보안팀은 통제와 증적을 얻고, 현업은 자동화 흐름을 유지할 수 있습니다.',
            '주요 구성은 Tool Call Firewall, Agent/Tool Registry, Execution Ontology, Policy Engine, Approval Workflow, Audit Timeline, Data Sensitivity Guard, Dashboard입니다.',
            '초기 고객은 AI 에이전트를 실제 업무 시스템에 연결하려는 금융, SOC/MSSP, 공공, 규제 산업, 대기업 DX 조직입니다. 이들은 자동화 수요와 감사·승인 요구를 동시에 갖고 있습니다.',
            '수익 모델은 엔터프라이즈 구독과 도입 전문서비스의 결합입니다. PoC로 시작해 정책 템플릿, 연동, 감사 리포트, 고급 컴플라이언스 기능으로 확장합니다.',
        ]),
        ('3. 시장 기회: Agent Execution Security의 등장', [
            'LLM 시장은 생성형 도구에서 에이전트 실행 계층으로 이동하고 있습니다. 에이전트가 외부 시스템을 호출할수록 통제의 중심은 프롬프트가 아니라 실행 요청으로 옮겨갑니다.',
            'MCP와 Tool Use는 연결을 쉽게 만들지만, 연결이 쉬워진 만큼 실행 리스크도 커집니다. 이메일 발송, DB 조회, 계정 변경, 배포, 파일 삭제는 모두 실제 피해를 만들 수 있는 행동입니다.',
            '기존 보안 예산은 IAM, PAM, DLP, CASB, SIEM/SOAR 같은 사람·네트워크·데이터 중심 통제에 집중되어 있습니다. 그러나 에이전트가 권한 안에서 내리는 실행 결정은 별도 계층이 필요합니다.',
            '우리는 이 범주를 Agent Execution Security 또는 Tool Call Governance로 정의합니다. 이는 기존 보안 제품을 대체하기보다, AI 에이전트 운영에 필요한 새로운 통제 지점을 추가합니다.',
            '시장 진입의 핵심은 거대한 수치보다 명확한 초기 문제입니다. 규제 산업과 보안운영 영역은 실행 통제와 감사 증적의 필요성을 가장 먼저 체감합니다.',
        ]),
        ('4. 문제 정의: Tool Call의 블랙박스', [
            '에이전트는 권한을 부여받는 순간부터 행동할 수 있는 존재가 됩니다. 문제는 그 권한 안에서 내려지는 판단이 항상 올바르지는 않다는 점입니다.',
            '프롬프트 인젝션, 과잉 권한, 잘못된 도구 선택, 민감정보 노출, 승인 없는 변경은 모두 에이전트 운영에서 발생할 수 있는 실행 리스크입니다.',
            '로그만으로는 충분하지 않습니다. 사고가 발생한 뒤 “무엇이 실행되었는가”를 보는 것보다, 실행되기 전에 “이 행동을 지금 허용해도 되는가”를 판단해야 합니다.',
            '기업이 필요한 것은 모델 보안만이 아닙니다. 모델이 생성한 행동이 실제 시스템에서 실행되는 순간을 통제하는 실행 보안입니다.',
        ]),
        ('5. 솔루션: Zero Trust Gateway for AI Agents', [
            'Cogito Zero는 에이전트와 도구/MCP/API 사이의 단일 통제 지점입니다. 에이전트가 호출하는 모든 고위험 Tool Call은 실행 전 정책 평가를 거칩니다.',
            '정책 엔진은 에이전트 신원, 사용자 권한, 도구 위험도, 데이터 민감도, 시간, 업무 목적, 과거 이력 같은 맥락을 함께 판단합니다.',
            '결과는 단순 허용·차단이 아닙니다. 안전한 호출은 Allow, 조건이 필요한 호출은 Conditional Allow, 고위험 호출은 Approval, 영향 확인이 필요한 호출은 Dry-run, 금지된 호출은 Block으로 처리합니다.',
            'Cogito Zero는 자동화를 막는 방화벽이 아니라, 자동화가 기업 통제 안에서 확장되도록 만드는 실행 거버넌스 레이어입니다.',
        ]),
        ('6. 제품 구성과 작동 방식', [
            'Tool Call Firewall은 에이전트의 호출을 실행 직전에 가로채고 정책 평가 대상으로 전환합니다. 이는 모든 실행 요청이 지나가는 공통 관문입니다.',
            'Agent/Tool Registry는 어떤 에이전트가 어떤 도구를 사용할 수 있는지 등록하고 신뢰 수준을 관리합니다. 조직은 흩어진 에이전트와 MCP 도구를 한 계층에서 볼 수 있습니다.',
            'Execution Ontology는 호출을 읽기, 쓰기, 전송, 삭제, 배포, 계정 변경, 외부 공유 같은 실행 의미로 분류합니다. 정책 엔진은 이 공통 언어 위에서 판단합니다.',
            'Approval Workflow는 위험한 요청을 사람 승인으로 전환합니다. Audit Timeline은 호출, 판단, 승인, 결과를 시간순으로 기록해 감사 가능한 증적을 남깁니다.',
            'Data Sensitivity Guard는 PII, 영업비밀, 규제 대상 데이터를 탐지하고 마스킹이나 전송 제한을 적용합니다. Dashboard는 보안팀과 컴플라이언스팀이 실행 리스크를 확인하는 관제 화면입니다.',
        ]),
        ('7. 정책 결정 체계', [
            '정책 결정은 다섯 가지로 단순화합니다. 이 구조는 보안팀, 현업, 감사 담당자가 같은 언어로 에이전트 실행을 이해하게 만듭니다.',
        ]),
        ('8. 대표 시나리오 3가지', [
            '첫 번째는 AI-SOC 안전 게이트입니다. 보안 에이전트가 야간에 관리자 계정 비활성화를 시도할 때, Cogito Zero는 고위험 대상과 업무 시간 맥락을 감지해 즉시 실행 대신 승인으로 전환합니다.',
            '두 번째는 고객정보 외부 전송 통제입니다. 고객지원 에이전트가 PII가 포함된 답변을 외부 채널로 보내려 할 때, 민감정보를 탐지해 마스킹하거나 담당자 승인을 요구합니다.',
            '세 번째는 MCP 도구 오염 차단입니다. 연결된 도구가 선언된 스키마와 다른 위험 동작을 요청하면, Cogito Zero는 스키마 불일치와 실행 의미를 비교해 호출을 차단합니다.',
            '이 시나리오의 공통 효과는 자동화 속도를 유지하면서도 되돌리기 어려운 실행을 사람과 정책의 통제 안에 둔다는 점입니다.',
        ]),
        ('9. 거버넌스와 컴플라이언스', [
            'AI 에이전트가 기업 시스템을 실행하는 순간, 감사 가능성은 선택이 아니라 운영 요구가 됩니다. 누가, 언제, 어떤 에이전트를 통해, 어떤 근거로 실행했는지 설명할 수 있어야 합니다.',
            'Cogito Zero의 정책, 승인, 감사 기능은 ISO/IEC 42001, IT 거버넌스, 개인정보 영향평가 같은 통제 요구에 대응할 수 있는 운영 증적을 제공합니다.',
            '정책은 보안팀만의 문서가 아니라 실행 시스템 안에서 작동해야 합니다. Approval Workflow와 Audit Timeline은 정책이 실제 행동에 적용되었음을 보여줍니다.',
            'Dry-run은 고위험 변경 전에 영향도를 확인하는 안전장치입니다. 이는 차단만으로는 해결하기 어려운 운영 연속성 문제를 줄입니다.',
        ]),
        ('10. 경쟁 환경과 포지셔닝', [
            'Cogito Zero는 IAM, PAM, API Gateway, DLP, LLM Security, Observability와 다른 계층을 담당합니다. 기존 제품은 각각 접근, 트래픽, 데이터, 모델 입출력, 사후 관찰에 강점이 있습니다.',
            '그러나 에이전트의 실행 행위를 정책 단위로 통제하고, 호출 전 판단·조건부 실행·승인·감사를 하나로 연결하는 계층은 아직 초기 시장입니다.',
            '따라서 Cogito Zero의 포지션은 “AI 에이전트 실행 권한을 관리하는 제로트러스트 게이트웨이”입니다. 기존 보안 체계와 연동하되 에이전트 시대의 빈 공간을 메우는 역할입니다.',
        ]),
        ('11. 비즈니스 모델', [
            '기본 수익 모델은 엔터프라이즈 구독입니다. 과금 기준은 보호 대상 에이전트 수, 연결 도구 수, Tool Call 규모, 감사 보존 기간, 고급 정책 기능 범위가 될 수 있습니다.',
            '초기 도입 단계에서는 PoC, 정책 설계, 기존 보안 도구 연동, 운영 교육을 전문서비스로 제공합니다. 이는 고객의 첫 적용 장벽을 낮추고 실제 사용 사례를 빠르게 확보하기 위한 구조입니다.',
            '상위 플랜은 고급 컴플라이언스 리포트, 전사 대시보드, 다중 승인, 장기 감사 보존, 전용 배포 옵션을 포함할 수 있습니다.',
        ]),
        ('12. GTM 및 초기 도입 전략', [
            '시장 진입은 고위험 워크플로우에서 시작합니다. 금융, SOC/MSSP, 공공, 대기업 DX 조직은 에이전트 자동화의 효익과 통제 요구를 동시에 갖고 있습니다.',
            '진입 방식은 Land and Expand입니다. 처음에는 CRM 변경, 고객정보 전송, 배포 자동화, 계정 변경 같은 제한된 워크플로우에 적용하고, 성공 사례를 바탕으로 팀·전사 단위로 확장합니다.',
            '개발자에게는 MCP/API Gateway로 쉽게 붙는 도구가 되어야 하고, 보안팀에게는 승인·감사·정책 리포트를 제공해야 합니다. 두 이해관계자를 동시에 설득하는 것이 핵심입니다.',
        ]),
        ('13. 로드맵과 팀', [
            'Phase 1은 핵심 게이트웨이입니다. Tool Call Firewall, Registry, 기본 Policy Engine, Audit Timeline으로 제한된 PoC를 검증합니다.',
            'Phase 2는 운영 확장입니다. Approval Workflow, Data Sensitivity Guard, Dashboard, 정책 템플릿, 주요 MCP/API 연동을 강화합니다.',
            'Phase 3은 플랫폼화입니다. 엔터프라이즈 권한 체계, 컴플라이언스 모듈, 파트너 생태계, 실행 리스크 분석을 통해 전사 표준 계층으로 확장합니다.',
            '팀 역량은 AI 에이전트 인프라, 보안 아키텍처, 엔터프라이즈 SaaS, 컴플라이언스 운영 경험을 중심으로 구성하는 것이 적합합니다.',
        ]),
        ('14. 재무 관점과 리스크', [
            '초기 재무 관점은 매출 과장보다 검증 가능한 지표에 집중해야 합니다. 핵심 지표는 protected tool calls, active agents, approval events, blocked risky actions, audit retention customers입니다.',
            '엔터프라이즈 보안 SaaS 특성상 매출은 파일럿에서 시작해 확장 계약으로 이어질 가능성이 큽니다. 따라서 1년차에는 PoC 성공률과 유료 전환율, 2년차 이후에는 계정 확장과 유지율이 중요합니다.',
            '주요 리스크는 에이전트 도입 속도 지연, 기존 보안 벤더의 기능 확장, 정책 설정 복잡성, 과도한 차단으로 인한 사용자 경험 저하입니다.',
            '대응책은 정책 템플릿, 드라이런 모드, 개발자 친화적 SDK/MCP 연동, 보안팀과 현업팀 모두가 이해할 수 있는 대시보드입니다.',
        ]),
        ('15. 맺음말', [
            'AI 에이전트가 기업 시스템을 실제로 실행하는 시대에는 새로운 신뢰 계층이 필요합니다. Cogito Zero는 그 계층을 제로트러스트 방식으로 제공합니다.',
            '목표는 에이전트를 막는 것이 아닙니다. 조직이 믿고 확장할 수 있도록, 위험한 실행은 통제하고 안전한 실행은 흐르게 만드는 것입니다.',
        ]),
    ]

    for title, ps in sections:
        h(doc, title)
        for ptext in ps:
            para(doc, ptext)
        if title == '2. 핵심 요약':
            table(doc, ['항목', 'Balanced 메시지'], [
                ['문제', '에이전트가 실제 도구를 실행하지만 실행 직전 통제와 감사 체계가 부족합니다.'],
                ['해법', '모든 Tool Call을 정책 기반으로 평가하고 승인·드라이런·차단으로 전환합니다.'],
                ['고객', '금융, SOC/MSSP, 공공, 규제 산업, 대기업 DX 조직부터 시작합니다.'],
                ['수익', '구독형 플랫폼 + PoC/정책 설계/연동 전문서비스로 진입합니다.'],
            ], [1.4, 5.0])
        if title == '7. 정책 결정 체계':
            table(doc, ['결정', '의미', '대표 상황'], [
                ['Allow', '정책상 안전한 실행을 즉시 통과', '읽기 전용 조회, 낮은 위험의 내부 검색'],
                ['Conditional Allow', '조건을 걸어 제한 실행', 'PII 마스킹 후 조회, 범위 제한 적용'],
                ['Approval', '사람 승인 후 실행', '계정 비활성화, 대량 삭제, 외부 전송'],
                ['Dry-run', '실행 전 영향 확인', '배포, 방화벽 규칙, 권한 변경 사전 검증'],
                ['Block', '정책상 금지된 실행 차단', '미승인 외부 공유, 스키마 불일치 도구 호출'],
            ], [1.2, 2.3, 3.0])
        if title == '10. 경쟁 환경과 포지셔닝':
            table(doc, ['카테고리', '주요 초점', 'Cogito Zero와의 차이'], [
                ['IAM/PAM', '사람과 계정 권한', '권한 안에서 발생하는 에이전트 실행 판단을 보완'],
                ['API Gateway', '인증, 라우팅, 트래픽', '에이전트 맥락과 데이터 민감도 기반 정책 판단 추가'],
                ['DLP/CASB', '데이터 이동과 SaaS 사용', 'Tool Call 단위의 실행 의도와 승인 흐름을 다룸'],
                ['LLM Security', '프롬프트/입출력 보안', '모델 답변 이후 실제 실행 계층을 통제'],
                ['Observability', '사후 관찰', '실행 전 판단과 승인·차단을 수행'],
            ], [1.5, 1.8, 3.2])

    doc.save(DOCX_OUT)


def add_bg(slide, color='bg'):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = prgb(color)


def tx(slide, x, y, w, h, text, size=16, color='charcoal', bold=False, align='left'):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for margin in ['margin_left', 'margin_right', 'margin_top', 'margin_bottom']:
        setattr(tf, margin, PptInches(0.03))
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    for r in p.runs:
        r.font.name = '맑은 고딕'
        r.font.size = PptPt(size)
        r.font.bold = bold
        r.font.color.rgb = prgb(color)
    return box


def ml(slide, x, y, w, h, lines, size=13, color='charcoal'):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = PptInches(0.03); tf.margin_right = PptInches(0.03); tf.margin_top = PptInches(0.02); tf.margin_bottom = PptInches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = PptPt(4)
        for r in p.runs:
            r.font.name = '맑은 고딕'
            r.font.size = PptPt(size)
            r.font.color.rgb = prgb(color)
    return box


def line(slide, x1, y1, x2, y2, color='line', width=.8):
    sh = slide.shapes.add_connector(1, PptInches(x1), PptInches(y1), PptInches(x2), PptInches(y2))
    sh.line.color.rgb = prgb(color)
    sh.line.width = PptPt(width)
    return sh


def header(slide, section, title, page):
    tx(slide, .65, .34, 2.6, .22, section.upper(), 8.2, 'teal', True)
    tx(slide, .65, .64, 9.4, .52, title, 25.5, 'charcoal', True)
    tx(slide, 11.55, 6.88, .85, .18, f'{page:02d} / 14', 7.5, 'slate', align='right')
    tx(slide, .65, 6.88, 3.4, .18, 'Cogito Zero · Balanced Business Plan', 7.5, 'slate')


def block(slide, x, y, w, label, title, body, accent=False):
    line(slide, x, y, x+w, y, 'teal' if accent else 'line', 1.1 if accent else .65)
    tx(slide, x, y+.15, w, .18, label, 7.8, 'slate', True)
    tx(slide, x, y+.43, w, .28, title, 12.8, 'charcoal', True)
    ml(slide, x, y+.78, w, .72, [body], 9.4, 'slate')


def pill(slide, x, y, w, text, fill='charcoal', color='white'):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(.38))
    sh.fill.solid(); sh.fill.fore_color.rgb = prgb(fill)
    sh.line.fill.background()
    tx(slide, x+.05, y+.08, w-.1, .18, text, 8.5, color, True, 'center')


def build_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    # 1 cover
    s = prs.slides.add_slide(blank); add_bg(s, 'charcoal')
    tx(s, .78, .62, 2.9, .24, 'BALANCED BUSINESS PLAN · 2026', 8.5, 'teal', True)
    tx(s, .78, 2.25, 6.2, .78, 'Cogito Zero', 46, 'white', True)
    ml(s, .82, 3.15, 6.6, .8, ['AI 에이전트와 도구/MCP/API 사이의', '제로트러스트 실행 통제 게이트웨이'], 17, 'white')
    line(s, .82, 5.7, 4.5, 5.7, 'teal', 1.0)
    tx(s, .82, 5.93, 5.2, .28, 'Governable · Auditable · Operationally Safe', 10.5, 'white')

    slides = [
        ('summary', '핵심 요약', [
            ('Problem', '실행 사각지대', '에이전트가 권한 안에서 고위험 Tool Call을 실행할 수 있습니다.'),
            ('Solution', '실행 직전 통제', '정책 평가, 승인, 드라이런, 차단, 감사를 하나의 흐름으로 연결합니다.'),
            ('Market', '규제 산업 우선', '금융, SOC, 공공, 대기업 DX처럼 증적 요구가 큰 곳부터 시작합니다.'),
            ('Business', 'PoC 후 확장', '구독형 플랫폼과 도입 전문서비스로 진입해 전사 표준으로 확장합니다.'),
        ], '자동화를 막지 않고, 위험한 실행만 통제합니다.'),
        ('market shift', 'AI 도입의 병목은 실행 신뢰성입니다', None, None),
        ('opportunity', 'Agent Execution Security의 등장', None, None),
        ('problem', 'Tool Call의 블랙박스', None, None),
        ('solution', 'Zero Trust Gateway for AI Agents', None, None),
        ('workflow', '작동 방식: Context → Policy → Decision → Audit', None, None),
        ('policy', '다섯 가지 결정으로 운영을 멈추지 않습니다', None, None),
        ('scenarios', '대표 시나리오 3가지', None, None),
        ('governance', '거버넌스와 컴플라이언스', None, None),
        ('positioning', '기존 보안의 틈을 메우는 계층', None, None),
        ('business', 'Business Model & GTM', None, None),
        ('financial', '재무 논리와 핵심 지표', None, None),
    ]
    # 2 summary
    s = prs.slides.add_slide(blank); add_bg(s); header(s, slides[0][0], slides[0][1], 2)
    for i, (lab, tit, body) in enumerate(slides[0][2]):
        block(s, .78 + i*3.0, 2.05, 2.35, lab, tit, body, i == 1)
    tx(s, .8, 5.72, 9.7, .42, slides[0][3], 22, 'charcoal', True)

    # 3 market shift
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'market shift', 'AI 도입의 병목은 실행 신뢰성입니다', 3)
    block(s, .9, 2.0, 4.2, 'FROM', 'AI Tools', '답변 생성, 요약, 초안 작성 중심. 위험은 주로 모델 입출력에 머물렀습니다.', False)
    block(s, 7.0, 2.0, 4.2, 'TO', 'Decision Systems', '에이전트가 도구를 호출하고 기업 시스템을 실제로 변경합니다.', True)
    line(s, 5.35, 3.0, 6.7, 3.0, 'teal', 1.2)
    tx(s, .9, 5.55, 9.8, .45, '질문은 “어떤 모델인가”에서 “이 실행을 지금 허용해도 되는가”로 이동합니다.', 20, 'charcoal', True)

    # 4 opportunity
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'opportunity', 'Agent Execution Security의 등장', 4)
    for i, (lab, tit, body) in enumerate([
        ('DRIVER 1', 'Automation pressure', '기업은 에이전트 자동화를 업무 시스템에 연결하려 합니다.'),
        ('DRIVER 2', 'Auditability', '실행 근거, 승인 이력, 결과 증적이 필요합니다.'),
        ('DRIVER 3', 'Regulation', 'AI 거버넌스와 개인정보 보호 요구가 강화됩니다.'),
    ]): block(s, .9+i*3.8, 2.05, 3.0, lab, tit, body, i==1)
    tx(s, .9, 5.65, 10.1, .42, 'Cogito Zero의 기회는 AI 도입 자체가 아니라, AI 실행을 통제하는 운영 계층입니다.', 18, 'charcoal', True)

    # 5 problem
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'problem', 'Tool Call의 블랙박스', 5)
    items = [('과잉 권한', '권한은 있지만 맥락상 위험한 실행'), ('민감정보', 'PII·영업비밀의 외부 전송'), ('오염된 도구', 'MCP 스키마와 실제 동작 불일치'), ('사후 로그 한계', '실행 전 승인·조건부 통제 부족')]
    for i, (t,b) in enumerate(items):
        x = .95 + (i%2)*5.25; y = 2.0 + (i//2)*1.45
        tx(s, x, y, 3.8, .3, t, 15, 'charcoal', True); ml(s, x, y+.42, 4.2, .45, [b], 11, 'slate'); line(s, x, y+.98, x+4.2, y+.98)
    tx(s, .95, 5.65, 9.0, .42, '로그만 남기는 것으로는 부족합니다. 실행되기 전에 판단해야 합니다.', 19, 'teal', True)

    # 6 solution
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'solution', 'Zero Trust Gateway for AI Agents', 6)
    tx(s, .95, 2.5, 2.0, .34, 'AI Agent', 18, 'charcoal', True, 'center')
    gate = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(4.55), PptInches(2.22), PptInches(3.45), PptInches(1.05)); gate.fill.solid(); gate.fill.fore_color.rgb = prgb('charcoal'); gate.line.fill.background()
    tx(s, 4.75, 2.55, 3.05, .28, 'Cogito Zero', 21, 'white', True, 'center')
    tx(s, 9.95, 2.5, 2.2, .34, 'Tools / MCP / API', 18, 'charcoal', True, 'center')
    line(s, 3.0, 2.75, 4.45, 2.75, 'teal', 1.2); line(s, 8.1, 2.75, 9.8, 2.75, 'teal', 1.2)
    ml(s, 1.05, 4.45, 10.6, .65, ['에이전트 신원, 도구 위험도, 데이터 민감도, 시간, 목적, 과거 이력을 함께 평가합니다.'], 16, 'charcoal')

    # 7 workflow
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'workflow', '작동 방식: Context → Policy → Decision → Audit', 7)
    steps = [('01', 'Context', '누가, 어떤 에이전트가, 어떤 도구로 무엇을 하려는가'), ('02', 'Policy', '조직 정책과 데이터 민감도, 도구 신뢰도를 평가'), ('03', 'Decision', 'Allow, Conditional, Approval, Dry-run, Block 중 결정'), ('04', 'Audit', '판단 근거와 승인, 실행 결과를 타임라인으로 기록')]
    for i,(n,t,b) in enumerate(steps):
        block(s, .8+i*3.0, 2.2, 2.35, n, t, b, i==2)
    tx(s, .82, 5.72, 8.8, .38, '복잡한 보안 판단을 하나의 실행 흐름으로 단순화합니다.', 18, 'charcoal', True)

    # 8 policy
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'policy', '다섯 가지 결정으로 운영을 멈추지 않습니다', 8)
    dec = [('Allow','즉시 허용'),('Conditional','조건부 실행'),('Approval','승인 요청'),('Dry-run','영향 확인'),('Block','차단')]
    for i,(a,b) in enumerate(dec):
        x=.8+i*2.42; tx(s,x,2.25,1.8,.25,f'{i+1:02d}',10,'teal',True); tx(s,x,2.72,1.8,.32,a,14,'charcoal',True); tx(s,x,3.12,1.8,.3,b,11,'slate'); line(s,x,3.72,x+1.75,3.72)
    tx(s,.85,5.58,10.1,.42,'목표는 모든 실행을 막는 것이 아니라, 위험도에 맞는 실행 방식을 선택하는 것입니다.',19,'charcoal',True)

    # 9 scenarios
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'scenarios', '대표 시나리오 3가지', 9)
    cases=[('AI-SOC','admin 계정 비활성화 요청','승인 전환','위협 대응과 업무 연속성 균형'),('PII Transfer','고객정보 외부 전송','마스킹·승인','개인정보 보호와 증적 확보'),('MCP Security','스키마 불일치 도구 호출','실행 차단','Tool Poisoning 실행 전 차단')]
    for i,(a,b,c,d) in enumerate(cases): block(s,.85+i*3.9,2.05,3.05,a,b,f'{c} → {d}',i==1)
    tx(s,.85,5.72,9.4,.38,'상황 → 판단 → 결과가 모두 감사 가능한 형태로 남습니다.',18,'charcoal',True)

    # 10 governance
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'governance', '거버넌스와 컴플라이언스', 10)
    for i,(t,b) in enumerate([('Policy','조직 정책을 실행 시스템 안에서 적용'),('Approval','고위험 행동을 사람 승인으로 전환'),('Logging','호출·판단·결과를 시간순으로 기록'),('Audit','누가·왜·무엇을 실행했는지 설명'),('Exception','예외와 dry-run으로 운영 연속성 보장')]):
        x=.9+(i%3)*3.8; y=1.95+(i//3)*1.55; block(s,x,y,3.0,f'{i+1:02d}',t,b,i==0)
    tx(s,.9,5.8,9.8,.35,'AI 사용을 막는 것이 아니라, 감사 가능한 방식으로 확장하게 만듭니다.',18,'teal',True)

    # 11 positioning
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'positioning', '기존 보안의 틈을 메우는 계층', 11)
    line(s,2.0,5.45,10.9,5.45,'slate',1); line(s,2.0,5.45,2.0,1.65,'slate',1)
    tx(s,1.35,1.35,1.2,.26,'실행',10,'slate',False,'center'); tx(s,10.2,5.65,1.7,.26,'AI Agent',10,'slate',False,'center')
    for text,x,y in [('IAM/PAM',3.0,4.7),('DLP/CASB',5.4,4.0),('LLM Security',7.4,3.25),('SIEM/SOAR',4.2,2.65)]: tx(s,x,y,1.9,.24,text,10.5,'slate',align='center')
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL,PptInches(8.65),PptInches(2.05),PptInches(.28),PptInches(.28)); dot.fill.solid(); dot.fill.fore_color.rgb=prgb('teal'); dot.line.fill.background()
    tx(s,8.95,2.02,2.1,.3,'Cogito Zero',15,'charcoal',True)
    tx(s,.9,5.95,9.5,.35,'모델 답변 이후, 실제 실행 직전에 정책 판단을 수행합니다.',17,'charcoal',True)

    # 12 business
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'business', 'Business Model & GTM', 12)
    for i,(lab,tit,body) in enumerate([('REVENUE','Platform subscription','에이전트 수, 보호 도구 수, Tool Call 규모, 감사 보존 기간 기반'),('SERVICES','PoC & integration','정책 설계, 보안 도구 연동, 운영 교육으로 초기 도입 장벽 완화'),('GTM','Land and expand','금융·SOC·공공의 고위험 워크플로우에서 시작해 전사 확장')]): block(s,.85+i*3.9,2.1,3.05,lab,tit,body,i==0)
    tx(s,.85,5.7,10.2,.35,'개발자에게는 쉽게 붙는 게이트웨이, 보안팀에게는 감사 가능한 통제 계층이 되어야 합니다.',17,'charcoal',True)

    # 13 financial
    s = prs.slides.add_slide(blank); add_bg(s); header(s, 'financial', '재무 논리와 핵심 지표', 13)
    for i,(phase,body) in enumerate([('Pilot','제한된 워크플로우에서 실행 통제 효과 검증'),('Expansion','유료 전환 후 에이전트·도구·정책 범위 확대'),('Platform','컴플라이언스 모듈과 장기 감사 보존으로 ARPA 상승')]): block(s,.9+i*3.8,2.0,3.0,phase.upper(),phase,body,i==1)
    ml(s,.9,4.95,8.8,.65,['핵심 지표: protected tool calls · active agents · approval events · blocked risky actions · audit retention customers'],13,'charcoal')
    tx(s,.9,5.85,9.8,.28,'초기에는 과장된 매출보다 PoC 성공률과 유료 전환율을 검증합니다.',16,'teal',True)

    # 14 roadmap closing
    s = prs.slides.add_slide(blank); add_bg(s,'charcoal')
    tx(s,.75,.58,2.1,.22,'ROADMAP',8.5,'teal',True)
    tx(s,.75,1.0,8.6,.48,'PoC에서 전사 실행 표준으로',28,'white',True)
    for i,(phase,body) in enumerate([('Foundation','Gateway, Registry, Policy Engine, Audit Timeline'),('Expansion','Approval, Data Guard, Dashboard, policy templates'),('Scale','Enterprise controls, integrations, decision intelligence')]):
        block(s,.85+i*3.9,2.25,3.05,phase.upper(),phase,body,i==0)
        # recolor text boxes from block are dark; overlay simple white headings for dark slide
        tx(s,.85+i*3.9,2.4,3.0,.24,phase.upper(),8,'teal',True)
        tx(s,.85+i*3.9,2.72,3.0,.3,phase,14,'white',True)
        ml(s,.85+i*3.9,3.1,3.0,.6,[body],9.3,'white')
    line(s,.85,5.35,4.8,5.35,'teal',1.1)
    tx(s,.85,5.62,9.8,.45,'Cogito Zero makes AI agent execution governable, auditable, and operationally safe.',20,'white',True)
    tx(s,11.55,6.88,.85,.18,'14 / 14',7.5,'white',align='right')

    prs.save(PPTX_OUT)


if __name__ == '__main__':
    build_docx()
    build_pptx()
    print(f'DOCX_OUT={DOCX_OUT}')
    print(f'PPTX_OUT={PPTX_OUT}')
